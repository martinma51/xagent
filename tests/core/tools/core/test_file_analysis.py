"""Tests for ``xagent.core.tools.core.file_analysis`` PowerPoint handling."""

from __future__ import annotations

import builtins
from pathlib import Path
from typing import Any, Iterable

import pytest
from pptx import Presentation
from pptx.util import Inches

from xagent.core.tools.core.file_analysis import (
    analyze_uploaded_file,
    collect_pptx_slide_blocks,
    iter_pptx_shapes,
)

# --------------------------------------------------------------------------- #
# Builders                                                                    #
# --------------------------------------------------------------------------- #


def _build_pptx(tmp_path: Path, builder) -> Path:
    """Run ``builder(prs)`` against a fresh Presentation, save it, return path."""
    prs = Presentation()
    builder(prs)
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))
    return pptx_path


def _add_title_body_slide(prs, title: str, body: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body


# --------------------------------------------------------------------------- #
# iter_pptx_shapes                                                            #
# --------------------------------------------------------------------------- #


def test_iter_pptx_shapes_descends_into_groups(tmp_path: Path) -> None:
    """Shapes inside a GroupShape are yielded as flat leaves."""

    def build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "top"
        tb1 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
        tb1.text_frame.text = "child A"
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
        tb2.text_frame.text = "child B"
        slide.shapes.add_group_shape([tb1, tb2])

    path = _build_pptx(tmp_path, build)
    prs = Presentation(str(path))
    texts = [
        getattr(shape, "text", "")
        for shape in iter_pptx_shapes(prs.slides[0].shapes)
        if getattr(shape, "text", "")
    ]
    assert "top" in texts
    assert "child A" in texts
    assert "child B" in texts


# --------------------------------------------------------------------------- #
# collect_pptx_slide_blocks                                                   #
# --------------------------------------------------------------------------- #


def test_collect_pptx_slide_blocks_text_tables_notes_groups(
    tmp_path: Path,
) -> None:
    """Helper extracts shape text, table cells, speaker notes, and group content."""

    def build(prs):
        _add_title_body_slide(prs, "Intro", "subtitle")
        prs.slides[0].notes_slide.notes_text_frame.text = "speaker notes"

        # Slide with a table
        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        s2.shapes.title.text = "Data"
        tbl = s2.shapes.add_table(
            2, 2, Inches(1), Inches(1), Inches(4), Inches(2)
        ).table
        tbl.cell(0, 0).text = "Name"
        tbl.cell(0, 1).text = "City"
        tbl.cell(1, 0).text = "Alice"
        tbl.cell(1, 1).text = "Shanghai"

        # Slide with a group
        s3 = prs.slides.add_slide(prs.slide_layouts[5])
        s3.shapes.title.text = "Grouped"
        tb1 = s3.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
        tb1.text_frame.text = "g-a"
        tb2 = s3.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
        tb2.text_frame.text = "g-b"
        s3.shapes.add_group_shape([tb1, tb2])

        # Empty slide
        prs.slides.add_slide(prs.slide_layouts[6])

    path = _build_pptx(tmp_path, build)
    blocks, stats = collect_pptx_slide_blocks(Presentation(str(path)))
    joined = "\n\n".join(blocks)

    assert "Intro" in joined
    assert "subtitle" in joined
    assert "[Notes] speaker notes" in joined
    assert "Alice" in joined
    assert "Shanghai" in joined
    assert "g-a" in joined
    assert "g-b" in joined
    # Empty slide must not produce a header
    assert "Slide 4:" not in joined

    assert stats["notes_count"] == 1
    assert stats["total_chars"] > 0
    assert stats["image_count"] == 0


def test_collect_pptx_slide_blocks_does_not_create_notes_slide(
    tmp_path: Path,
) -> None:
    """The analyzer must not mutate the deck by accessing slide.notes_slide."""

    def build(prs):
        _add_title_body_slide(prs, "no notes", "body")

    path = _build_pptx(tmp_path, build)
    collect_pptx_slide_blocks(Presentation(str(path)))

    # Reload from disk and verify no notes slide was added
    reloaded = Presentation(str(path))
    assert reloaded.slides[0].has_notes_slide is False


# --------------------------------------------------------------------------- #
# analyze_uploaded_file - .pptx branch                                        #
# --------------------------------------------------------------------------- #


def test_analyze_pptx_returns_preview_and_summary(tmp_path: Path) -> None:
    def build(prs):
        _add_title_body_slide(prs, "Hello", "World")

    path = _build_pptx(tmp_path, build)
    result = analyze_uploaded_file(path.name, uploads_dir=str(path.parent))

    assert result.file_type == ".pptx"
    assert "Hello" in result.content_preview
    assert "World" in result.content_preview
    assert "PowerPoint presentation" in result.structure_summary
    assert any("slide(s) extracted" in i for i in result.key_insights)


def test_analyze_pptx_corrupt_file_returns_error_result(tmp_path: Path) -> None:
    bad = tmp_path / "corrupt.pptx"
    bad.write_bytes(b"not actually a pptx zip")

    result = analyze_uploaded_file(bad.name, uploads_dir=str(tmp_path))

    assert result.file_type == ".pptx"
    assert result.content_preview == "[Failed to open presentation]"
    assert "Invalid or corrupt" in result.structure_summary
    assert any("python-pptx could not open" in i for i in result.key_insights)


# --------------------------------------------------------------------------- #
# analyze_uploaded_file - .ppt legacy fallback                                #
# --------------------------------------------------------------------------- #


def test_analyze_legacy_ppt_returns_binary_fallback(tmp_path: Path) -> None:
    """Legacy .ppt files take a direct binary-fallback path; no UTF-8 read."""
    ppt = tmp_path / "legacy.ppt"
    # CFB header bytes — content doesn't have to be a real .ppt
    ppt.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1binary blob")

    result = analyze_uploaded_file(ppt.name, uploads_dir=str(tmp_path))

    assert result.file_type == ".ppt"
    assert result.content_preview == "[Binary or non-text file]"
    assert "Legacy PowerPoint file" in result.structure_summary
    assert result.key_insights[0].startswith("Legacy .ppt format detected")


# --------------------------------------------------------------------------- #
# analyze_uploaded_file - ImportError fallback                                #
# --------------------------------------------------------------------------- #


def test_analyze_pptx_without_python_pptx_falls_back(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """If python-pptx is unavailable at runtime, the analyzer degrades."""

    def build(prs):
        _add_title_body_slide(prs, "T", "B")

    path = _build_pptx(tmp_path, build)

    real_import = builtins.__import__

    def fake_import(
        name: str,
        globals_: Any = None,
        locals_: Any = None,
        fromlist: Iterable[str] = (),
        level: int = 0,
    ):
        if name == "pptx" or name.startswith("pptx."):
            raise ImportError(name)
        return real_import(name, globals_, locals_, fromlist, level)

    monkeypatch.setattr(builtins, "__import__", fake_import)

    result = analyze_uploaded_file(path.name, uploads_dir=str(path.parent))

    assert result.file_type == ".pptx"
    assert any("python-pptx not installed" in i for i in result.key_insights)
