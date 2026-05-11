"""File analysis tools for agent"""

import json
from pathlib import Path
from typing import Any, Dict, Iterator, List, Tuple

from pydantic import BaseModel


class FileAnalysisResult(BaseModel):
    """File analysis result model"""

    filename: str
    file_type: str
    size: int
    content_preview: str
    structure_summary: str
    key_insights: List[str]
    suggested_actions: List[str]


def analyze_uploaded_file(
    filename: str, uploads_dir: str = "/uploads"
) -> FileAnalysisResult:
    """
    Analyze uploaded file content

    Args:
        filename: File name
        uploads_dir: Upload directory path

    Returns:
        File analysis result

    Raises:
        FileNotFoundError: File doesn't exist
        ValueError: File type not supported
    """
    # Build file path
    file_path = Path(uploads_dir) / filename

    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {filename}")

    # Get file info
    stat = file_path.stat()
    file_size = stat.st_size
    file_ext = file_path.suffix.lower()

    # Analyze based on file type
    if file_ext == ".txt":
        return _analyze_text_file(file_path, filename, file_size)
    elif file_ext == ".json":
        return _analyze_json_file(file_path, filename, file_size)
    elif file_ext == ".csv":
        return _analyze_csv_file(file_path, filename, file_size)
    elif file_ext == ".py":
        return _analyze_python_file(file_path, filename, file_size)
    elif file_ext == ".md":
        return _analyze_markdown_file(file_path, filename, file_size)
    elif file_ext in (".pptx", ".ppt"):
        return _analyze_pptx_file(file_path, filename, file_size, file_ext)
    else:
        return _analyze_generic_file(file_path, filename, file_size, file_ext)


def _analyze_text_file(
    file_path: Path, filename: str, file_size: int
) -> FileAnalysisResult:
    """Analyze text file"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        # Generate content preview
        content_preview = content[:500] + "..." if len(content) > 500 else content

        # Analyze structure
        lines = content.split("\n")
        word_count = len(content.split())

        structure_summary = f"Text file with {len(lines)} lines, {word_count} words"

        # Generate insights
        insights = []
        if len(lines) > 100:
            insights.append("File is quite long - consider summarizing key sections")
        if word_count > 1000:
            insights.append(
                "Document contains substantial content - may require structured analysis"
            )

        # Suggested actions
        actions = [
            "Read complete content for detailed analysis",
            "Extract key information or summaries",
            "Search for specific keywords or patterns",
        ]

        return FileAnalysisResult(
            filename=filename,
            file_type="text",
            size=file_size,
            content_preview=content_preview,
            structure_summary=structure_summary,
            key_insights=insights,
            suggested_actions=actions,
        )

    except Exception as e:
        raise ValueError(f"Failed to analyze text file: {e}")


def _analyze_json_file(
    file_path: Path, filename: str, file_size: int
) -> FileAnalysisResult:
    """Analyze JSON file"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        # Generate content preview
        content_preview = json.dumps(data, indent=2, ensure_ascii=False)[:500] + "..."

        # Analyze structure
        if isinstance(data, dict):
            keys = list(data.keys())
            structure_summary = f"JSON object with {len(keys)} keys: {keys[:5]}{'...' if len(keys) > 5 else ''}"

            # Check nested structure
            max_depth = _get_json_depth(data)
            if max_depth > 3:
                structure_summary += f", max depth: {max_depth}"

        elif isinstance(data, list):
            structure_summary = f"JSON array with {len(data)} items"
            if data and isinstance(data[0], dict):
                structure_summary += f" (objects with {len(data[0])} fields each)"
        else:
            structure_summary = f"JSON {type(data).__name__}"

        # Generate insights
        insights = []
        if isinstance(data, dict) and len(data) > 20:
            insights.append(
                "Large JSON structure - consider data processing or transformation"
            )
        if isinstance(data, list) and len(data) > 100:
            insights.append("Large dataset - consider filtering or aggregation")

        # Suggested actions
        actions = [
            "Parse and extract specific fields",
            "Transform or restructure data",
            "Validate data integrity",
            "Generate reports or summaries",
        ]

        return FileAnalysisResult(
            filename=filename,
            file_type="json",
            size=file_size,
            content_preview=content_preview,
            structure_summary=structure_summary,
            key_insights=insights,
            suggested_actions=actions,
        )

    except Exception as e:
        raise ValueError(f"Failed to analyze JSON file: {e}")


def _analyze_csv_file(
    file_path: Path, filename: str, file_size: int
) -> FileAnalysisResult:
    """Analyze CSV file"""
    try:
        import csv

        with open(file_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)

        if not rows:
            raise ValueError("CSV file is empty")

        # Get basic information
        header = rows[0]
        total_rows = len(rows) - 1
        total_cols = len(header)

        # Generate content preview
        preview_rows = rows[:6]  # First 5 rows of data including header
        content_preview = "\n".join([",".join(row) for row in preview_rows])
        if len(rows) > 6:
            content_preview += "\n..."

        structure_summary = f"CSV with {total_rows} data rows, {total_cols} columns: {header[:5]}{'...' if len(header) > 5 else ''}"

        # Generate insights
        insights = []
        if total_rows > 1000:
            insights.append("Large dataset - consider sampling or filtering")
        if total_cols > 20:
            insights.append(
                "Many columns - consider feature selection or dimensionality reduction"
            )

        # Suggested actions
        actions = [
            "Analyze data patterns and statistics",
            "Clean and preprocess data",
            "Extract specific columns or rows",
            "Generate visualizations or reports",
        ]

        return FileAnalysisResult(
            filename=filename,
            file_type="csv",
            size=file_size,
            content_preview=content_preview,
            structure_summary=structure_summary,
            key_insights=insights,
            suggested_actions=actions,
        )

    except Exception as e:
        raise ValueError(f"Failed to analyze CSV file: {e}")


def _analyze_python_file(
    file_path: Path, filename: str, file_size: int
) -> FileAnalysisResult:
    """Analyze Python file"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        # Analyze Python code structure
        lines = content.split("\n")

        # Collect basic statistics
        import_count = content.count("import ")
        from_count = content.count("from ")
        class_count = content.count("class ")
        def_count = content.count("def ")

        # Generate content preview
        content_preview = content[:500] + "..." if len(content) > 500 else content

        structure_summary = f"Python file with {len(lines)} lines, {class_count} classes, {def_count} functions"

        # Generate insights
        insights = []
        if len(lines) > 500:
            insights.append("Large Python file - consider modularization")
        if import_count + from_count > 10:
            insights.append("Many dependencies - check for unused imports")

        # Suggested actions
        actions = [
            "Execute code to test functionality",
            "Analyze code quality and structure",
            "Extract documentation or comments",
            "Generate code summaries or refactoring suggestions",
        ]

        return FileAnalysisResult(
            filename=filename,
            file_type="python",
            size=file_size,
            content_preview=content_preview,
            structure_summary=structure_summary,
            key_insights=insights,
            suggested_actions=actions,
        )

    except Exception as e:
        raise ValueError(f"Failed to analyze Python file: {e}")


def _analyze_markdown_file(
    file_path: Path, filename: str, file_size: int
) -> FileAnalysisResult:
    """Analyze Markdown file"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        # Analyze Markdown structure
        lines = content.split("\n")

        # Count headers
        headers = [line for line in lines if line.startswith("#")]

        # Count code blocks
        code_blocks = content.count("```")

        # Generate content preview
        content_preview = content[:500] + "..." if len(content) > 500 else content

        structure_summary = f"Markdown with {len(headers)} headers, {code_blocks // 2} code blocks, {len(lines)} lines"

        # Generate insights
        insights = []
        if len(headers) > 10:
            insights.append("Well-structured document with clear hierarchy")
        if code_blocks > 4:
            insights.append("Contains multiple code examples or snippets")

        # Suggested actions
        actions = [
            "Extract structured information",
            "Convert to other formats",
            "Generate table of contents",
            "Summarize key sections",
        ]

        return FileAnalysisResult(
            filename=filename,
            file_type="markdown",
            size=file_size,
            content_preview=content_preview,
            structure_summary=structure_summary,
            key_insights=insights,
            suggested_actions=actions,
        )

    except Exception as e:
        raise ValueError(f"Failed to analyze Markdown file: {e}")


def iter_pptx_shapes(shapes: Any) -> Iterator[Any]:
    """Yield individual shapes, recursing into ``GroupShape`` containers.

    Top-level iteration over ``slide.shapes`` misses content nested inside
    grouped shapes (a common PowerPoint pattern). This generator descends into
    any group it encounters so callers see a flat stream of leaf shapes.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_pptx_shapes(shape.shapes)
        else:
            yield shape


def collect_pptx_slide_blocks(prs: Any) -> Tuple[List[str], Dict[str, int]]:
    """Walk a python-pptx ``Presentation`` and return per-slide text blocks.

    Returns:
        A pair ``(blocks, stats)`` where ``blocks`` is a list of formatted
        per-slide strings (e.g. ``"Slide 1:\nTitle\n..."``). Slides with no
        text / notes / table content are omitted. ``stats`` is a dict with
        ``total_shapes``, ``total_chars``, ``notes_count``, ``image_count``.

    Callers are expected to have verified that ``python-pptx`` is importable
    before passing a Presentation in.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    blocks: list[str] = []
    total_shapes = 0
    total_chars = 0
    notes_count = 0
    image_count = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = []
        for shape in iter_pptx_shapes(slide.shapes):
            total_shapes += 1
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            shape_text = getattr(shape, "text", None)
            if shape_text:
                txt = str(shape_text).strip()
                if txt:
                    slide_lines.append(txt)
                    total_chars += len(txt)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        txt = cell.text_frame.text.strip()
                        if txt:
                            slide_lines.append(txt)
                            total_chars += len(txt)
        # ``slide.notes_slide`` creates a notes slide on access; guard with
        # ``has_notes_slide`` so this analyzer stays read-only.
        notes = ""
        if getattr(slide, "has_notes_slide", False):
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            notes = getattr(notes_frame, "text", "").strip() if notes_frame else ""
        if notes:
            notes_count += 1
            slide_lines.append(f"[Notes] {notes}")
            total_chars += len(notes)
        if slide_lines:
            blocks.append(f"Slide {slide_num}:\n" + "\n".join(slide_lines))

    stats = {
        "total_shapes": total_shapes,
        "total_chars": total_chars,
        "notes_count": notes_count,
        "image_count": image_count,
    }
    return blocks, stats


def _analyze_pptx_file(
    file_path: Path, filename: str, file_size: int, file_ext: str
) -> FileAnalysisResult:
    """Analyze a PowerPoint presentation (.pptx / .ppt) file.

    Uses python-pptx to extract slide titles, body text and speaker notes so the
    agent receives a usable text representation instead of binary blob bytes.
    Legacy ``.ppt`` (binary) files are not supported by python-pptx; they fall
    back to the generic analyzer with a hint to convert.
    """
    if file_ext == ".ppt":
        # python-pptx only supports the OOXML .pptx format. Return a binary
        # fallback directly instead of letting _analyze_generic_file attempt a
        # UTF-8 read on the binary CFB stream.
        return FileAnalysisResult(
            filename=filename,
            file_type=file_ext,
            size=file_size,
            content_preview="[Binary or non-text file]",
            structure_summary=(
                f"Legacy PowerPoint file ({file_ext}) with {file_size} bytes"
            ),
            key_insights=[
                "Legacy .ppt format detected; convert to .pptx for full text extraction.",
                "File appears to be binary or encoded",
            ],
            suggested_actions=[
                "Use appropriate binary file tools",
                "Convert to .pptx format",
            ],
        )

    try:
        from pptx import Presentation
    except ImportError:
        result = _analyze_generic_file(file_path, filename, file_size, file_ext)
        result.key_insights.insert(
            0,
            "python-pptx not installed; install xagent[document-processing] to parse .pptx.",
        )
        return result

    try:
        prs = Presentation(str(file_path))
    except Exception as exc:  # noqa: BLE001
        return FileAnalysisResult(
            filename=filename,
            file_type=file_ext,
            size=file_size,
            content_preview="[Failed to open presentation]",
            structure_summary=f"Invalid or corrupt .pptx file ({file_size} bytes)",
            key_insights=[f"python-pptx could not open file: {exc}"],
            suggested_actions=[
                "Verify the file is a valid .pptx",
                "Re-export from PowerPoint or Keynote",
            ],
        )

    slide_blocks, stats = collect_pptx_slide_blocks(prs)
    total_shapes = stats["total_shapes"]
    total_chars = stats["total_chars"]
    notes_count = stats["notes_count"]
    image_count = stats["image_count"]

    full_text = "\n\n".join(slide_blocks)
    content_preview = (full_text[:500] + "...") if len(full_text) > 500 else full_text

    structure_summary = (
        f"PowerPoint presentation with {len(prs.slides)} slide(s), "
        f"{total_shapes} shape(s), {image_count} picture(s), "
        f"{notes_count} slide(s) with speaker notes ({total_chars} text chars)"
    )

    insights = [
        f"{len(prs.slides)} slide(s) extracted via python-pptx",
        f"{total_chars} characters of text content available",
    ]
    if notes_count:
        insights.append(f"{notes_count} slide(s) contain speaker notes")
    if image_count:
        insights.append(
            f"{image_count} embedded picture(s) - consider OCR if text is in images"
        )
    if total_chars == 0:
        insights.append("No extractable text - slides may consist of images only")

    actions = [
        "Use the extracted slide text as agent context",
        "Pass to document_parser tool with parser_name='unstructured' for richer chunks",
        "Render to PDF/HTML for visual preview if needed",
    ]
    if image_count:
        actions.append("Run image OCR on slides whose text lives inside images")

    return FileAnalysisResult(
        filename=filename,
        file_type=file_ext,
        size=file_size,
        content_preview=content_preview,
        structure_summary=structure_summary,
        key_insights=insights,
        suggested_actions=actions,
    )


def _analyze_generic_file(
    file_path: Path, filename: str, file_size: int, file_ext: str
) -> FileAnalysisResult:
    """Analyze generic file"""
    try:
        # Try to read first few lines as preview
        with open(file_path, "r", encoding="utf-8") as f:
            preview_lines = []
            for i, line in enumerate(f):
                if i >= 10:  # Only read first 10 lines
                    break
                preview_lines.append(line.rstrip())

        content_preview = "\n".join(preview_lines)
        if len(preview_lines) >= 10:
            content_preview += "\n..."

        structure_summary = f"Generic file ({file_ext}) with {file_size} bytes"

        # Generate insights
        insights = [
            f"File type {file_ext} may require specialized handling",
            "Consider using appropriate tools for this file format",
        ]

        # Suggested actions
        actions = [
            "Use file-type specific tools if available",
            "Convert to more common format if possible",
            "Extract text content for analysis",
        ]

        return FileAnalysisResult(
            filename=filename,
            file_type=file_ext,
            size=file_size,
            content_preview=content_preview,
            structure_summary=structure_summary,
            key_insights=insights,
            suggested_actions=actions,
        )

    except Exception:
        # If unable to read as text, provide basic information
        return FileAnalysisResult(
            filename=filename,
            file_type=file_ext,
            size=file_size,
            content_preview="[Binary or non-text file]",
            structure_summary=f"Binary file ({file_ext}) with {file_size} bytes",
            key_insights=["File appears to be binary or encoded"],
            suggested_actions=["Use appropriate binary file tools"],
        )


def _get_json_depth(obj: Any, current_depth: int = 0) -> int:
    """Get maximum depth of JSON data"""
    if isinstance(obj, dict):
        if not obj:
            return current_depth
        return max(_get_json_depth(v, current_depth + 1) for v in obj.values())
    elif isinstance(obj, list):
        if not obj:
            return current_depth
        return max(_get_json_depth(item, current_depth + 1) for item in obj)
    else:
        return current_depth


def list_uploaded_files(uploads_dir: str = "/uploads") -> Dict[str, Any]:
    """
    List files in upload directory

    Args:
        uploads_dir: Upload directory path

    Returns:
        File list information
    """
    upload_path = Path(uploads_dir)

    if not upload_path.exists():
        return {
            "files": [],
            "total_count": 0,
            "message": "Uploads directory does not exist",
        }

    files = []
    for file_path in upload_path.iterdir():
        if file_path.is_file():
            stat = file_path.stat()
            files.append(
                {
                    "filename": file_path.name,
                    "size": stat.st_size,
                    "modified_time": stat.st_mtime,
                    "extension": file_path.suffix.lower(),
                    "file_path": str(file_path),
                }
            )

    return {
        "files": files,
        "total_count": len(files),
        "uploads_directory": str(upload_path.absolute()),
    }


def get_file_context(filename: str, uploads_dir: str = "/uploads") -> Dict[str, Any]:
    """
    Get complete context information for a file

    Args:
        filename: File name
        uploads_dir: Upload directory path

    Returns:
        File context information
    """
    try:
        # Analyze file
        analysis = analyze_uploaded_file(filename, uploads_dir)

        # Get file list
        files_info = list_uploaded_files(uploads_dir)

        return {
            "analysis": analysis.model_dump(),
            "available_files": files_info,
            "file_access_info": {
                "filename": filename,
                "full_path": str(Path(uploads_dir) / filename),
                "uploads_directory": uploads_dir,
                "accessible": True,
            },
        }

    except Exception as e:
        return {
            "error": str(e),
            "filename": filename,
            "accessible": False,
            "available_files": list_uploaded_files(uploads_dir),
        }
